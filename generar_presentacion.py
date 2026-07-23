"""
Generador de guías operativas
=============================

Rellena una plantilla .pptx con marcadores {{llave}} y elimina los bloques de
las variables que no traen datos, reacomodando el diseño para que no queden
huecos.

No requiere renombrar figuras en PowerPoint: los bloques visuales se detectan
por su posición vertical, usando bandas que se calculan a partir de la
separación real entre los elementos de cada lista.

    python3 generador.py
"""

import os
import re
import copy
from statistics import median
from pptx import Presentation

NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
PATRON_LLAVE = re.compile(r"\{\{([^{}]+)\}\}")
PATRON_MARCADOR = re.compile(r"\{\{.*?\}\}")

# Una figura más alta que esta fracción de la diapositiva se considera
# estructural (barra lateral, fondo) y nunca se elimina ni se usa como
# contenedor de sección.
FRACCION_ESTRUCTURAL = 0.5

# Margen extra sobre la altura de la banda al decidir qué figuras pertenecen
# a la fila de un elemento.
FACTOR_BANDA = 1.5


# ---------------------------------------------------------------------------
# TEXTO: reemplazo conservando el formato original
# ---------------------------------------------------------------------------

def _run_representativo(parrafo):
    """Run cuyo formato debe conservarse.

    Los marcadores suelen venir partidos en varios runs ('{{' + 'norma_1' +
    '}}') y el primero a veces es solo la llave de apertura. Se toma el run
    con más caracteres reales, que es el que lleva el formato del contenido.
    """
    if not parrafo.runs:
        return None
    return max(parrafo.runs, key=lambda r: len(re.sub(r"[^\w]", "", r.text)))


def reemplazar_texto_manteniendo_formato(parrafo, texto_nuevo):
    runs = parrafo.runs
    if not runs:
        parrafo.text = texto_nuevo
        return

    modelo = _run_representativo(parrafo)
    destino = runs[0]

    if modelo is not None and modelo is not destino:
        rPr_modelo = modelo._r.find(f"{NS_A}rPr")
        if rPr_modelo is not None:
            rPr_destino = destino._r.find(f"{NS_A}rPr")
            if rPr_destino is not None:
                destino._r.remove(rPr_destino)
            destino._r.insert(0, copy.deepcopy(rPr_modelo))

    destino.text = texto_nuevo
    for run in runs[1:]:
        run._r.getparent().remove(run._r)


def procesar_parrafo(parrafo, datos):
    texto = parrafo.text
    if "{{" not in texto:
        return
    for llave, valor in datos.items():
        marcador = "{{" + llave + "}}"
        if marcador in texto:
            texto = texto.replace(marcador, str(valor))
    texto = PATRON_MARCADOR.sub("", texto)
    reemplazar_texto_manteniendo_formato(parrafo, texto.strip())


def activar_ajuste_al_texto(shape):
    """Deja la caja en 'ajustar forma al texto' para que su alto se recalcule
    según el contenido cuando PowerPoint abra el archivo."""
    try:
        bodyPr = shape.text_frame._txBody.find(f"{NS_A}bodyPr")
        if bodyPr is None:
            return
        for etiqueta in ("noAutofit", "normAutofit"):
            viejo = bodyPr.find(f"{NS_A}{etiqueta}")
            if viejo is not None:
                bodyPr.remove(viejo)
        if bodyPr.find(f"{NS_A}spAutoFit") is None:
            bodyPr.append(bodyPr.makeelement(f"{NS_A}spAutoFit", {}))
    except Exception:
        pass


# ---------------------------------------------------------------------------
# GEOMETRÍA
# ---------------------------------------------------------------------------

def centro_y(shape):
    if shape.top is None or shape.height is None:
        return None
    return shape.top + shape.height // 2


def es_estructural(shape, alto_slide):
    """Barra lateral, fondo: figuras que atraviesan la diapositiva."""
    return shape.height is not None and shape.height > alto_slide * FRACCION_ESTRUCTURAL


def eliminar_figura(shape):
    padre = shape._element.getparent()
    if padre is not None:
        padre.remove(shape._element)


def partes_de_llave(llave):
    """'norma_2' -> ('norma', 2); 'objetivo' -> ('objetivo', 0)."""
    m = re.match(r"^(.*?)_(\d+)$", llave)
    return (m.group(1), int(m.group(2))) if m else (llave, 0)


def cajas_con_marcador(slide):
    """[(llave, shape)] de las cajas de texto que contienen un {{marcador}}."""
    salida = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for llave in PATRON_LLAVE.findall(shape.text_frame.text):
            salida.append((llave.strip(), shape))
    return salida


def agrupar_en_listas(cajas):
    """{prefijo: [(indice, llave, shape, centro), ...]} ordenado por posición."""
    listas = {}
    for llave, shape in cajas:
        c = centro_y(shape)
        if c is None:
            continue
        prefijo, indice = partes_de_llave(llave)
        listas.setdefault(prefijo, []).append((indice, llave, shape, c))
    for items in listas.values():
        items.sort(key=lambda x: x[3])
    return listas


def bandas_de_lista(items):
    """Banda vertical de cada elemento, según sus vecinos de la misma lista."""
    centros = [c for _, _, _, c in items]
    if len(centros) > 1:
        seps = [b - a for a, b in zip(centros, centros[1:])]
        medio = median(seps) / 2
    else:
        medio = items[0][2].height * 3

    bandas = []
    for i, c in enumerate(centros):
        inicio = (centros[i - 1] + c) / 2 if i > 0 else c - medio
        fin = (c + centros[i + 1]) / 2 if i < len(centros) - 1 else c + medio
        bandas.append((inicio, fin))
    return bandas


def figuras_en_banda(slide, banda, alto_slide):
    """Figuras que pertenecen a la fila visual de un elemento.

    Se excluyen las estructurales (barra lateral) y las demasiado altas para
    la banda (contenedores de sección), de modo que sólo se toman el recuadro
    del elemento, su viñeta o círculo, el número y los espacios de imagen.
    """
    inicio, fin = banda
    alto_max = (fin - inicio) * FACTOR_BANDA
    seleccion = []
    for shape in slide.shapes:
        if es_estructural(shape, alto_slide):
            continue
        c = centro_y(shape)
        if c is None or not (inicio <= c <= fin):
            continue
        if shape.height > alto_max:
            continue
        seleccion.append(shape)
    return seleccion


def contenedor_de_lista(slide, items, alto_slide):
    """Recuadro que engloba a toda una lista (ej. la tarjeta de 'Normas').

    Se busca la figura sin texto más pequeña que contenga verticalmente a
    todos los elementos de la lista. Devuelve None si no existe.
    """
    arriba = min(sh.top for _, _, sh, _ in items)
    abajo = max(sh.top + sh.height for _, _, sh, _ in items)

    candidatos = []
    for shape in slide.shapes:
        if es_estructural(shape, alto_slide):
            continue
        if shape.has_text_frame and shape.text_frame.text.strip():
            continue
        if shape.top is None or shape.height is None:
            continue
        if shape.top <= arriba and shape.top + shape.height >= abajo:
            candidatos.append(shape)

    return min(candidatos, key=lambda s: s.height) if candidatos else None


def figuras_dentro_de(slide, contenedor, alto_slide):
    """Todo lo que vive dentro de un contenedor (incluido él mismo)."""
    arriba = contenedor.top
    abajo = contenedor.top + contenedor.height
    dentro = []
    for shape in slide.shapes:
        if es_estructural(shape, alto_slide):
            continue
        c = centro_y(shape)
        if c is not None and arriba <= c <= abajo:
            dentro.append(shape)
    return dentro


# ---------------------------------------------------------------------------
# FLUJO PRINCIPAL
# ---------------------------------------------------------------------------

def tiene_datos(llave, datos):
    return llave in datos and str(datos[llave]).strip() != ""


def eliminar_diapositiva(prs, index):
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def procesar_slide(slide, datos, alto_slide, reacomodar=True):
    cajas = cajas_con_marcador(slide)
    if not cajas:
        return False

    listas = agrupar_en_listas(cajas)

    # Si ninguna variable de la diapositiva trae datos, se elimina entera.
    if not any(tiene_datos(llave, datos) for llave, _ in cajas):
        return True

    # (bottom_y, alto_liberado): espacio que libera cada sección al encoger.
    liberaciones = []

    for _, items in listas.items():
        bandas = bandas_de_lista(items)
        ranuras = [sh.top for _, _, sh, _ in items]
        activos = [i for i, (_, llave, _, _) in enumerate(items)
                   if tiene_datos(llave, datos)]

        if not activos:
            # Lista completamente vacía: se elimina su recuadro y su contenido
            # (título de la sección incluido).
            contenedor = contenedor_de_lista(slide, items, alto_slide)
            if contenedor is not None:
                base = contenedor.top + contenedor.height
                for shape in figuras_dentro_de(slide, contenedor, alto_slide):
                    eliminar_figura(shape)
                liberaciones.append((base, contenedor.height))
            else:
                for i in range(len(items)):
                    for shape in figuras_en_banda(slide, bandas[i], alto_slide):
                        eliminar_figura(shape)
            continue

        if len(activos) == len(items):
            continue  # nada que eliminar en esta lista

        contenedor = contenedor_de_lista(slide, items, alto_slide)

        # 1. Se guardan las figuras de cada elemento ANTES de borrar nada,
        #    para que las bandas se calculen sobre el diseño completo.
        bloques = {i: figuras_en_banda(slide, bandas[i], alto_slide)
                   for i in range(len(items))}

        # 2. Se eliminan los bloques de los elementos sin datos: recuadro,
        #    viñeta o círculo, número, texto y espacios de imagen.
        for i in range(len(items)):
            if i not in activos:
                for shape in bloques[i]:
                    eliminar_figura(shape)

        # 3. Los elementos que quedan suben a ocupar las primeras ranuras.
        for destino, origen in enumerate(activos):
            ancla = items[origen][2]
            delta = ranuras[destino] - ancla.top
            if delta == 0:
                continue
            for shape in bloques[origen]:
                if shape.top is not None:
                    shape.top = shape.top + delta

        # 4. El recuadro de la sección se encoge por el espacio liberado.
        liberado = ranuras[-1] - ranuras[len(activos) - 1]
        if liberado > 0 and contenedor is not None:
            base = contenedor.top + contenedor.height
            contenedor.height = max(contenedor.height - liberado, 1)
            liberaciones.append((base, liberado))

    # 5. Todo lo que está debajo de una sección encogida sube esa distancia.
    if reacomodar and liberaciones:
        for shape in slide.shapes:
            if shape.top is None or es_estructural(shape, alto_slide):
                continue
            desplazamiento = sum(alto for base, alto in liberaciones
                                 if base <= shape.top)
            if desplazamiento:
                shape.top = shape.top - desplazamiento

    # 6. Se rellenan los marcadores restantes conservando el formato.
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tenia_marcador = "{{" in shape.text_frame.text
        for parrafo in shape.text_frame.paragraphs:
            procesar_parrafo(parrafo, datos)
        if tenia_marcador:
            activar_ajuste_al_texto(shape)

    return False


def generar_guia_operativa(ruta_plantilla, ruta_salida, datos, reacomodar=True):
    if not datos or not isinstance(datos, dict):
        print("ADVERTENCIA: no se recibieron datos. Operación cancelada.")
        return

    prs = Presentation(ruta_plantilla)
    alto_slide = prs.slide_height
    a_eliminar = []

    for idx, slide in enumerate(prs.slides):
        if procesar_slide(slide, datos, alto_slide, reacomodar):
            a_eliminar.append(idx)

    for index in reversed(a_eliminar):
        eliminar_diapositiva(prs, index)

    prs.save(ruta_salida)
    print(f"Guía operativa generada en: {ruta_salida}")
    if a_eliminar:
        print(f"Diapositivas eliminadas por quedar sin contenido: {len(a_eliminar)}")


# ---------------------------------------------------------------------------
# PRUEBA
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    base = os.path.dirname(os.path.abspath(__file__))
    plantilla = os.path.join(base, "plantilla.pptx")
    salida = os.path.join(base, "Guia_Final.pptx")

    datos_prueba = {
        "Nombre_proceso": "Lavado de manos",
        "Fecha_vigencia": "22/02/2030",
        "codigo": "GO-RES-001",
        "revisión": "1",
        "tiempo_ejecucion": "5 minutos",
        "objetivo": "Establecer el método y técnica correcta para el lavado de manos.",

        "norma_1": "Cumplir el procedimiento antes de iniciar actividades.",
        # norma_2 ausente: su viñeta debe desaparecer y norma_3 subir.
        "norma_3": "Asegurar disponibilidad de servilletas industriales.",

        "responsable_1": "Gerente de restaurante",
        "responsable_2": "Manipulador de alimentos",

        "material_1": "Jabón bactericida.",
        "material-2": "Servilletas industriales.",

        "equipo_1": "Malla o gorro.",
        "equipo_2": "Zapatos de seguridad.",

        # Lista de tips completamente vacía: su recuadro debe eliminarse.

        "titulo_conocimiento_1": "Contaminación",
        "texto_conocimiento_1": "Alteración dañina de un medio por un agente ajeno.",
        "titulo_conocimiento_2": "Contaminación cruzada",
        "texto_conocimiento_2": "Transferencia de microbios de alimentos crudos a listos.",
        "titulo_conocimiento_3": "Microorganismo",
        "texto_conocimiento_3": "Seres vivos visibles solo al microscopio.",

        "contexto_operativo": "Garantiza condiciones higiénicas mínimas antes de cualquier tarea.",

        "descripcion_paso_1": "Mojar las manos con agua limpia.",
        "descripcion_paso_2": "Presionar dos (2) veces el dispensador de jabón bactericida.",
        "descripcion_paso_3": "Aplicar el jabón en manos, antebrazos y entre los dedos.",
        "descripcion_paso_4": "Enjuagar completamente el jabón con agua limpia.",
        "descripcion_paso_5": "Secar manos y antebrazos con servilletas industriales.",
        # pasos 6 y 7 ausentes: sus bloques completos deben desaparecer.
    }

    if os.path.exists(plantilla):
        generar_guia_operativa(plantilla, salida, datos_prueba)
    else:
        print(f"No se encontró {plantilla}")