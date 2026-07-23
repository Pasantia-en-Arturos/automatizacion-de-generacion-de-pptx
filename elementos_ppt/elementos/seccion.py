"""
Recuadro de sección
====================

El contenedor blanco de bordes redondeados que agrupa visualmente cada
sección (Objetivo, Normas, Responsables...).
"""

from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from . import estilo


def crear_recuadro_seccion(slide, left, top, width, height,
                            radio_esquina=0.14):
    """Rectángulo redondeado blanco con borde gris claro, como el que
    envuelve cada sección en la plantilla original."""
    forma = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
    )
    forma.fill.solid()
    forma.fill.fore_color.rgb = estilo.BLANCO
    forma.line.color.rgb = estilo.GRIS_BORDE
    forma.line.width = Pt(3)
    forma.adjustments[0] = radio_esquina
    forma.shadow.inherit = False
    return forma


def crear_punto_vineta(slide, left, top, diametro=Emu(220200)):
    """El punto naranja de una lista con viñetas."""
    punto = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(left), Emu(top), diametro, diametro
    )
    punto.fill.solid()
    punto.fill.fore_color.rgb = estilo.NARANJA
    punto.line.fill.background()
    return punto
