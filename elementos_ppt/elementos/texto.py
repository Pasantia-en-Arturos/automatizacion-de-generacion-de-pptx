"""
Cuadros de texto
================

Funciones para crear cajas de texto desde cero con python-pptx, usando la
tipografía y color de la marca (ver estilo.py). Cada función devuelve el
shape creado, por si se necesita seguir ajustándolo después.
"""

from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from . import estilo


def crear_cuadro_titulo(slide, texto, left, top, width, height,
                         tamano=estilo.TAMANO_TITULO_SECCION,
                         color=estilo.NARANJA):
    """Un título de sección, ej. 'Objetivo:', 'Normas:'."""
    caja = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = caja.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = texto
    run.font.name = estilo.FUENTE_TITULO
    run.font.size = Pt(tamano)
    run.font.bold = True
    run.font.color.rgb = color
    return caja


def crear_cuadro_texto(slide, texto, left, top, width, height,
                        tamano=estilo.TAMANO_CUERPO,
                        color=estilo.GRIS_CUERPO,
                        marcador=None):
    """Un cuadro de texto de cuerpo. Si se pasa `marcador` (ej. 'objetivo'),
    el texto real se reemplaza por '{{marcador}}' — útil para ir armando la
    plantilla con sus etiquetas ya puestas desde el propio código."""
    caja = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = caja.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{{{{{marcador}}}}}" if marcador else texto
    run.font.name = estilo.FUENTE_CUERPO
    run.font.size = Pt(tamano)
    run.font.color.rgb = color
    return caja
