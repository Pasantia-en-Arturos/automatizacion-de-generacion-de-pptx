"""
Demostración
============

Arma una diapositiva completa desde cero, combinando únicamente los módulos
de elementos/ — sin partir de ninguna plantilla ya diseñada.

    python3 demo.py
"""

from pptx import Presentation
from pptx.util import Emu

from elementos.texto import crear_cuadro_titulo, crear_cuadro_texto
from elementos.seccion import crear_recuadro_seccion
from elementos.lista import crear_item_lista

prs = Presentation()
prs.slide_width = Emu(19202400)   # mismas proporciones que la plantilla real
prs.slide_height = Emu(27432000)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout en blanco

# --- Sección "Objetivo" ---
crear_recuadro_seccion(slide, left=1060350, top=5749187,
                        width=17033700, height=2200000)
crear_cuadro_titulo(slide, "Objetivo:", left=1593318, top=5974554,
                     width=8000000, height=923400)
crear_cuadro_texto(slide, "", left=1593325, top=6827361,
                    width=15000000, height=664767, marcador="objetivo")

# --- Sección "Normas" con 2 elementos de lista ---
crear_recuadro_seccion(slide, left=1060350, top=8895962,
                        width=17033700, height=3800000)
crear_cuadro_titulo(slide, "Normas:", left=1593318, top=9035464,
                     width=8000000, height=923400)

crear_item_lista(slide, left=1930627, top=10551057,
                  width=15000000, height=664767, marcador="norma_1")
crear_item_lista(slide, left=1930627, top=11400000,
                  width=15000000, height=664767, marcador="norma_2")

prs.save("demo_generado.pptx")
print("Diapositiva generada en: demo_generado.pptx")
