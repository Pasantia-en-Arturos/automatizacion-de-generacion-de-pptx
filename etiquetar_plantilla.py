"""
Etiquetador de plantillas
=========================

Asigna a cada figura auxiliar (punto de viñeta, círculo del paso, recuadro
contenedor, espacio para imagen) un nombre con la convención:

    var:<llave>

de modo que el generador sepa a qué variable pertenece cada elemento visual
sin depender de su posición en la diapositiva.

Se ejecuta UNA SOLA VEZ por plantilla. Después conviene abrir el resultado en
PowerPoint (Inicio > Organizar > Panel de selección) y revisar los nombres.

    python3 etiquetar_plantilla.py plantilla.pptx plantilla_etiquetada.pptx
"""

import re
import sys
from statistics import median
from pptx import Presentation

PATRON_MARCADOR = re.compile(r"\{\{([\w\-áéíóúñÁÉÍÓÚÑ]+)\}\}")

# Textos de relleno del diseño que SÍ deben eliminarse junto con su bloque.
ETIQUETAS_DE_RELLENO = ("ESPACIO PARA", "IMAGEN/GRÁFICO")

# Longitud máxima de texto para considerar una figura como "decorativa"
# (ej. el número "6" del paso). Evita etiquetar texto fijo del diseño.
MAX_TEXTO_DECORATIVO = 3


def prefijo_de_lista(llave: str):
    """'norma_2' -> ('norma', 2). Devuelve None si la llave no es de lista."""
    m = re.match(r"^(.*?)_(\d+)$", llave)
    return (m.group(1), int(m.group(2))) if m else None


def centro_vertical(shape):
    if shape.top is None or shape.height is None:
        return None
    return shape.top + shape.height // 2


def es_decorativa(shape):
    """True si la figura puede pertenecer al bloque de una variable."""
    if not shape.has_text_frame:
        return True
    texto = shape.text_frame.text.strip()
    if not texto:
        return True
    if PATRON_MARCADOR.search(texto):
        return False  # es la caja de otra variable
    if any(e in texto.upper() for e in (t.upper() for t in ETIQUETAS_DE_RELLENO)):
        return True
    return len(texto) <= MAX_TEXTO_DECORATIVO


def cajas_de_marcador(slide):
    """[(llave, shape, centro)] de las cajas que contienen un {{marcador}}."""
    salida = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for llave in PATRON_MARCADOR.findall(shape.text_frame.text):
            c = centro_vertical(shape)
            if c is not None:
                salida.append((llave, shape, c))
    return salida


def bandas_por_lista(cajas):
    """Calcula la banda vertical de cada variable de lista.

    La banda se define por los puntos medios respecto a los elementos vecinos
    de la MISMA lista (norma_1, norma_2, norma_3...), de forma que se
    autocalibra al diseño de cada plantilla.
    """
    listas = {}
    for llave, shape, centro in cajas:
        info = prefijo_de_lista(llave)
        if info:
            listas.setdefault(info[0], []).append((info[1], llave, shape, centro))

    bandas = {}
    for prefijo, items in listas.items():
        items.sort(key=lambda x: x[0])
        centros = [c for _, _, _, c in items]
        if len(centros) > 1:
            seps = [b - a for a, b in zip(centros, centros[1:])]
            medio = median(seps) / 2
        else:
            medio = items[0][2].height * 3

        for i, (_, llave, _, centro) in enumerate(items):
            inicio = (centros[i - 1] + centro) / 2 if i > 0 else centro - medio
            fin = (centro + centros[i + 1]) / 2 if i < len(centros) - 1 else centro + medio
            bandas[llave] = (inicio, fin)
    return bandas


def etiquetar(ruta_entrada: str, ruta_salida: str):
    prs = Presentation(ruta_entrada)
    total = 0

    for slide in prs.slides:
        cajas = cajas_de_marcador(slide)
        bandas = bandas_por_lista(cajas)
        cajas_marcador = {id(sh) for _, sh, _ in cajas}

        # 1. La propia caja del marcador lleva el nombre de su variable.
        for llave, shape, _ in cajas:
            shape.name = f"var:{llave}"
            total += 1

        # 2. Las figuras auxiliares dentro de la banda de cada variable.
        for llave, (inicio, fin) in bandas.items():
            alto_max = (fin - inicio) * 1.5
            for shape in slide.shapes:
                if id(shape) in cajas_marcador:
                    continue
                if shape.name.startswith("var:"):
                    continue  # ya asignada a otra variable
                c = centro_vertical(shape)
                if c is None or not (inicio <= c <= fin):
                    continue
                if shape.height > alto_max:
                    continue  # barra lateral / contenedor de sección
                if not es_decorativa(shape):
                    continue
                shape.name = f"var:{llave}"
                total += 1

    prs.save(ruta_salida)
    print(f"Figuras etiquetadas: {total}")
    print(f"Plantilla etiquetada guardada en: {ruta_salida}")
    print("Revisa los nombres en PowerPoint: Inicio > Organizar > Panel de selección")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Uso: python3 etiquetar_plantilla.py entrada.pptx salida.pptx")
        sys.exit(1)
    etiquetar(sys.argv[1], sys.argv[2])